"""
Analisador de Templates IB - Acessa o Google Drive, baixa os .pptx de exemplo
e extrai todas as regras de formatação (fonts, cores, layouts, posições).

Gera:
  1. Um JSON detalhado com as especificações visuais
  2. Um resumo em texto para usar como system prompt no agente

Pré-requisitos:
  1. Criar um projeto no Google Cloud Console (console.cloud.google.com)
  2. Ativar a API Google Drive
  3. Criar credenciais OAuth 2.0 (tipo "Desktop App")
  4. Em "Authorized redirect URIs" do seu OAuth client, ADICIONAR:
       http://localhost:8080/
     (sem isso ocorre Erro 400: redirect_uri_mismatch)
  5. Baixar o arquivo credentials.json para a pasta scripts/
  6. pip install google-api-python-client google-auth-httplib2 google-auth-oauthlib python-pptx
"""

import os
import io
import json
import sys
from collections import Counter, defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OAUTH_LOCAL_PORT = 8080  # Porta fixa para o redirect_uri (registre http://localhost:8080/ no Google Cloud Console)
SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
CREDENTIALS_FILE = os.path.join(SCRIPT_DIR, "credentials.json")
TOKEN_FILE = os.path.join(SCRIPT_DIR, "token.json")
OUTPUT_DIR = os.path.join(SCRIPT_DIR, "..", "pptx-service", "templates")
ANALYSIS_OUTPUT = os.path.join(SCRIPT_DIR, "..", "template_analysis.json")
PROMPT_OUTPUT = os.path.join(SCRIPT_DIR, "..", "template_prompt.txt")


# ================================================================
# Google Drive - Autenticação e Download
# ================================================================

def authenticate_gdrive():
    """Autentica com Google Drive e retorna o serviço."""
    from google.auth.transport.requests import Request
    from google.oauth2.credentials import Credentials
    from google_auth_oauthlib.flow import InstalledAppFlow
    from googleapiclient.discovery import build

    creds = None
    if os.path.exists(TOKEN_FILE):
        creds = Credentials.from_authorized_user_file(TOKEN_FILE, SCOPES)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            if not os.path.exists(CREDENTIALS_FILE):
                print(f"ERRO: Arquivo credentials.json não encontrado em {CREDENTIALS_FILE}")
                print("Siga as instruções no topo deste script para configurar.")
                sys.exit(1)
            flow = InstalledAppFlow.from_client_secrets_file(CREDENTIALS_FILE, SCOPES)
            print(f"\n  Abrindo navegador para login. Redirect URI esperado: http://localhost:{OAUTH_LOCAL_PORT}/")
            print(f"  Se der erro redirect_uri_mismatch, adicione esse URI em Google Cloud Console > Credentials > seu OAuth client > Authorized redirect URIs.\n")
            creds = flow.run_local_server(port=OAUTH_LOCAL_PORT)
        with open(TOKEN_FILE, "w") as token:
            token.write(creds.to_json())

    return build("drive", "v3", credentials=creds)


def find_pptx_files(service, folder_id=None):
    """Busca arquivos .pptx no Drive (opcionalmente dentro de uma pasta)."""
    query_parts = [
        "mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'"
    ]
    if folder_id:
        query_parts.append(f"'{folder_id}' in parents")
    query_parts.append("trashed=false")
    query = " and ".join(query_parts)

    results = service.files().list(
        q=query,
        pageSize=20,
        fields="files(id, name, size, modifiedTime)",
        orderBy="modifiedTime desc",
    ).execute()

    return results.get("files", [])


def download_pptx(service, file_id, file_name, dest_dir):
    """Baixa um arquivo .pptx do Drive."""
    from googleapiclient.http import MediaIoBaseDownload

    os.makedirs(dest_dir, exist_ok=True)
    dest_path = os.path.join(dest_dir, file_name)

    request = service.files().get_media(fileId=file_id)
    buffer = io.BytesIO()
    downloader = MediaIoBaseDownload(buffer, request)

    done = False
    while not done:
        _, done = downloader.next_chunk()

    with open(dest_path, "wb") as f:
        f.write(buffer.getvalue())

    print(f"  Baixado: {file_name} ({len(buffer.getvalue()) / 1024:.0f} KB)")
    return dest_path


# ================================================================
# Análise de Templates PPTX
# ================================================================

def emu_to_inches(emu):
    """Converte EMU para polegadas."""
    if emu is None:
        return None
    return round(emu / 914400, 2)


def emu_to_pt(emu):
    """Converte EMU para pontos."""
    if emu is None:
        return None
    return round(emu / 12700, 1)


def rgb_to_hex(rgb_color):
    """Converte RGBColor para string hex."""
    if rgb_color is None:
        return None
    try:
        return f"#{rgb_color}"
    except Exception:
        return None


def extract_font_info(font):
    """Extrai informações de uma fonte."""
    info = {}
    try:
        if font.name:
            info["name"] = font.name
    except Exception:
        pass
    try:
        if font.size:
            info["size_pt"] = emu_to_pt(font.size)
    except Exception:
        pass
    try:
        info["bold"] = font.bold
    except Exception:
        pass
    try:
        info["italic"] = font.italic
    except Exception:
        pass
    try:
        if font.color and font.color.rgb:
            info["color"] = rgb_to_hex(font.color.rgb)
    except Exception:
        pass
    return info


def extract_paragraph_info(paragraph):
    """Extrai informações de um parágrafo."""
    info = {}
    try:
        if paragraph.alignment is not None:
            align_map = {
                PP_ALIGN.LEFT: "left",
                PP_ALIGN.CENTER: "center",
                PP_ALIGN.RIGHT: "right",
                PP_ALIGN.JUSTIFY: "justify",
            }
            info["alignment"] = align_map.get(paragraph.alignment, str(paragraph.alignment))
    except Exception:
        pass

    try:
        if paragraph.space_before:
            info["space_before_pt"] = emu_to_pt(paragraph.space_before)
    except Exception:
        pass
    try:
        if paragraph.space_after:
            info["space_after_pt"] = emu_to_pt(paragraph.space_after)
    except Exception:
        pass
    try:
        if paragraph.line_spacing:
            info["line_spacing"] = paragraph.line_spacing
    except Exception:
        pass

    runs = []
    for run in paragraph.runs:
        run_info = {"text_preview": run.text[:80] if run.text else ""}
        run_info["font"] = extract_font_info(run.font)
        runs.append(run_info)
    if runs:
        info["runs"] = runs

    return info


def extract_shape_info(shape):
    """Extrai informações de um shape."""
    info = {
        "shape_type": str(shape.shape_type) if shape.shape_type else None,
        "name": shape.name,
        "position": {
            "left_inches": emu_to_inches(shape.left),
            "top_inches": emu_to_inches(shape.top),
            "width_inches": emu_to_inches(shape.width),
            "height_inches": emu_to_inches(shape.height),
        },
    }

    try:
        if shape.fill.type is not None:
            info["fill_type"] = str(shape.fill.type)
            try:
                if shape.fill.fore_color and shape.fill.fore_color.rgb:
                    info["fill_color"] = rgb_to_hex(shape.fill.fore_color.rgb)
            except Exception:
                pass
    except Exception:
        pass

    if shape.has_text_frame:
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            p_info = extract_paragraph_info(para)
            if p_info:
                paragraphs.append(p_info)
        if paragraphs:
            info["paragraphs"] = paragraphs

    if shape.has_table:
        table = shape.table
        info["table"] = {
            "rows": len(table.rows),
            "cols": len(table.columns),
        }
        if len(table.rows) > 0 and len(table.columns) > 0:
            header_fonts = []
            for ci in range(len(table.columns)):
                cell = table.cell(0, ci)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        header_fonts.append(extract_font_info(r.font))
            if header_fonts:
                info["table"]["header_fonts"] = header_fonts
            try:
                cell = table.cell(0, 0)
                if cell.fill.type is not None:
                    info["table"]["header_fill"] = rgb_to_hex(cell.fill.fore_color.rgb)
            except Exception:
                pass

    if shape.has_chart:
        chart = shape.chart
        info["chart"] = {
            "chart_type": str(chart.chart_type),
            "has_legend": chart.has_legend,
        }

    return info


def analyze_slide(slide, slide_idx):
    """Analisa um slide completo."""
    slide_info = {
        "slide_number": slide_idx + 1,
        "layout_name": slide.slide_layout.name if slide.slide_layout else None,
    }

    try:
        bg = slide.background
        if bg.fill.type is not None:
            slide_info["background_type"] = str(bg.fill.type)
            try:
                if bg.fill.fore_color and bg.fill.fore_color.rgb:
                    slide_info["background_color"] = rgb_to_hex(bg.fill.fore_color.rgb)
            except Exception:
                pass
    except Exception:
        pass

    shapes = []
    for shape in slide.shapes:
        s_info = extract_shape_info(shape)
        shapes.append(s_info)

    slide_info["shapes"] = shapes
    slide_info["shape_count"] = len(shapes)

    return slide_info


def analyze_presentation(filepath):
    """Analisa uma apresentação inteira."""
    prs = Presentation(filepath)
    filename = os.path.basename(filepath)

    analysis = {
        "filename": filename,
        "slide_width_inches": emu_to_inches(prs.slide_width),
        "slide_height_inches": emu_to_inches(prs.slide_height),
        "total_slides": len(prs.slides),
        "slide_layouts_available": [],
        "slides": [],
    }

    for layout in prs.slide_layouts:
        layout_info = {
            "name": layout.name,
            "placeholders": [],
        }
        for ph in layout.placeholders:
            ph_info = {
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "name": ph.name,
                "position": {
                    "left_inches": emu_to_inches(ph.left),
                    "top_inches": emu_to_inches(ph.top),
                    "width_inches": emu_to_inches(ph.width),
                    "height_inches": emu_to_inches(ph.height),
                },
            }
            layout_info["placeholders"].append(ph_info)
        analysis["slide_layouts_available"].append(layout_info)

    for i, slide in enumerate(prs.slides):
        analysis["slides"].append(analyze_slide(slide, i))

    return analysis


# ================================================================
# Consolidação e Geração de Resumo
# ================================================================

def consolidate_analyses(analyses):
    """Consolida as análises de múltiplos arquivos em regras unificadas."""
    all_fonts = Counter()
    all_font_sizes = Counter()
    all_colors = Counter()
    all_bg_colors = Counter()
    all_fill_colors = Counter()
    title_fonts = []
    body_fonts = []

    for analysis in analyses:
        for slide in analysis["slides"]:
            if "background_color" in slide:
                all_bg_colors[slide["background_color"]] += 1

            for shape in slide.get("shapes", []):
                if "fill_color" in shape:
                    all_fill_colors[shape["fill_color"]] += 1

                for para in shape.get("paragraphs", []):
                    for run in para.get("runs", []):
                        font = run.get("font", {})
                        if "name" in font:
                            all_fonts[font["name"]] += 1
                        if "size_pt" in font:
                            all_font_sizes[font["size_pt"]] += 1
                        if "color" in font and font["color"]:
                            all_colors[font["color"]] += 1

                        is_title = (
                            "title" in shape.get("name", "").lower()
                            or (font.get("size_pt") and font["size_pt"] >= 20)
                        )
                        if is_title:
                            title_fonts.append(font)
                        else:
                            body_fonts.append(font)

    def most_common_value(items, key):
        vals = [item.get(key) for item in items if item.get(key) is not None]
        if not vals:
            return None
        return Counter(vals).most_common(1)[0][0]

    consolidated = {
        "dimensions": {
            "width_inches": analyses[0]["slide_width_inches"],
            "height_inches": analyses[0]["slide_height_inches"],
        },
        "fonts_used": dict(all_fonts.most_common(10)),
        "font_sizes_used": dict(all_font_sizes.most_common(15)),
        "text_colors_used": dict(all_colors.most_common(10)),
        "background_colors": dict(all_bg_colors.most_common(5)),
        "fill_colors": dict(all_fill_colors.most_common(10)),
        "title_style": {
            "font": most_common_value(title_fonts, "name"),
            "size_pt": most_common_value(title_fonts, "size_pt"),
            "bold": most_common_value(title_fonts, "bold"),
            "color": most_common_value(title_fonts, "color"),
        },
        "body_style": {
            "font": most_common_value(body_fonts, "name"),
            "size_pt": most_common_value(body_fonts, "size_pt"),
            "bold": most_common_value(body_fonts, "bold"),
            "color": most_common_value(body_fonts, "color"),
        },
        "total_files_analyzed": len(analyses),
        "total_slides_analyzed": sum(a["total_slides"] for a in analyses),
    }

    return consolidated


def generate_prompt_text(consolidated, analyses):
    """Gera texto formatado com regras de estilo para uso no system prompt."""
    c = consolidated
    dim = c["dimensions"]
    title = c["title_style"]
    body = c["body_style"]

    fonts = ", ".join(c["fonts_used"].keys()) if c["fonts_used"] else "Calibri"
    bg_colors = ", ".join(c["background_colors"].keys()) if c["background_colors"] else "Não identificado"
    fill_colors = list(c["fill_colors"].keys())[:6] if c["fill_colors"] else []
    text_colors = list(c["text_colors_used"].keys())[:6] if c["text_colors_used"] else []

    layouts_seen = set()
    for a in analyses:
        for s in a["slides"]:
            if s.get("layout_name"):
                layouts_seen.add(s["layout_name"])

    text = f"""## REGRAS DE FORMATAÇÃO IB (Extraídas de {c['total_files_analyzed']} apresentações de exemplo, {c['total_slides_analyzed']} slides)

### Dimensões
- Slide: {dim['width_inches']}" x {dim['height_inches']}" ({"widescreen 16:9" if dim['width_inches'] > 12 else "4:3"})

### Tipografia
- Fontes usadas: {fonts}
- Títulos: {title.get('font', 'N/A')} {title.get('size_pt', 'N/A')}pt{' bold' if title.get('bold') else ''}, cor {title.get('color', 'N/A')}
- Corpo: {body.get('font', 'N/A')} {body.get('size_pt', 'N/A')}pt{' bold' if body.get('bold') else ''}, cor {body.get('color', 'N/A')}

### Paleta de Cores
- Backgrounds: {bg_colors}
- Cores de preenchimento (shapes): {', '.join(fill_colors) if fill_colors else 'N/A'}
- Cores de texto: {', '.join(text_colors) if text_colors else 'N/A'}

### Tamanhos de Fonte Encontrados
- {', '.join(f'{s}pt' for s in sorted(c['font_sizes_used'].keys(), reverse=True))}

### Layouts de Slide Encontrados
- {', '.join(sorted(layouts_seen)) if layouts_seen else 'Blank / Custom'}

### Diretrizes Gerais
- Mantenha consistência visual com as cores e fontes acima
- Títulos sempre na fonte e tamanho especificados
- Corpo do texto no tamanho e cor especificados
- Use a paleta de cores para gráficos, tabelas e destaques
"""
    return text


# ================================================================
# Main
# ================================================================

def main():
    print("=" * 60)
    print("  Analisador de Templates IB - Google Drive")
    print("=" * 60)

    # Passo 1: Autenticar e buscar arquivos
    print("\n[1/4] Autenticando com Google Drive...")
    service = authenticate_gdrive()

    # Perguntar se quer buscar em pasta específica
    folder_id = None
    folder_input = input(
        "\nDigite o ID da pasta do Google Drive (ou Enter para buscar em todo o Drive): "
    ).strip()
    if folder_input:
        folder_id = folder_input

    print("\n[2/4] Buscando arquivos .pptx...")
    files = find_pptx_files(service, folder_id)

    if not files:
        print("Nenhum arquivo .pptx encontrado!")
        sys.exit(1)

    print(f"\nEncontrados {len(files)} arquivo(s):")
    for i, f in enumerate(files):
        size = int(f.get("size", 0)) / 1024
        print(f"  {i + 1}. {f['name']} ({size:.0f} KB) - Modificado: {f['modifiedTime'][:10]}")

    # Passo 2: Baixar arquivos
    print(f"\n[3/4] Baixando para {OUTPUT_DIR}...")
    downloaded = []
    for f in files:
        path = download_pptx(service, f["id"], f["name"], OUTPUT_DIR)
        downloaded.append(path)

    # Passo 3: Analisar
    print("\n[4/4] Analisando templates...")
    analyses = []
    for path in downloaded:
        print(f"  Analisando: {os.path.basename(path)}")
        analysis = analyze_presentation(path)
        analyses.append(analysis)
        print(f"    -> {analysis['total_slides']} slides, "
              f"{analysis['slide_width_inches']}\" x {analysis['slide_height_inches']}\"")

    # Passo 4: Consolidar e salvar
    consolidated = consolidate_analyses(analyses)

    full_output = {
        "consolidated": consolidated,
        "individual_analyses": analyses,
    }

    os.makedirs(os.path.dirname(ANALYSIS_OUTPUT), exist_ok=True)
    with open(ANALYSIS_OUTPUT, "w", encoding="utf-8") as f:
        json.dump(full_output, f, indent=2, ensure_ascii=False)
    print(f"\nAnálise detalhada salva em: {ANALYSIS_OUTPUT}")

    prompt_text = generate_prompt_text(consolidated, analyses)
    with open(PROMPT_OUTPUT, "w", encoding="utf-8") as f:
        f.write(prompt_text)
    print(f"Regras para system prompt salvas em: {PROMPT_OUTPUT}")

    # Mostrar resumo
    print("\n" + "=" * 60)
    print("  RESUMO DA ANÁLISE")
    print("=" * 60)
    print(prompt_text)
    print("=" * 60)
    print("Próximos passos:")
    print("  1. Revise o arquivo template_analysis.json para detalhes completos")
    print("  2. Copie o conteúdo de template_prompt.txt para o system prompt do agente N8N")
    print("  3. Ajuste as cores e fontes no generator.py se necessário")


if __name__ == "__main__":
    main()
